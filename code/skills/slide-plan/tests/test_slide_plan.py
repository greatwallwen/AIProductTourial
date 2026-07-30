from __future__ import annotations

import importlib.util
import json
import os
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
BUILD_SCRIPT = ROOT / "scripts" / "build_deck.mjs"
AUDIT_SCRIPT = ROOT / "scripts" / "audit_pptx.py"
NODE = Path(os.environ["CODEX_BUNDLED_NODE"]) if os.environ.get("CODEX_BUNDLED_NODE") else None
SETUP = Path(os.environ["CODEX_ARTIFACT_SETUP"]) if os.environ.get("CODEX_ARTIFACT_SETUP") else None
ARTIFACT_RUNTIME_AVAILABLE = bool(NODE and NODE.is_file() and SETUP and SETUP.is_file())


def load_module(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    sys.modules[name] = module
    spec.loader.exec_module(module)
    return module


AUDIT = load_module("audit_pptx", AUDIT_SCRIPT)


MARKDOWN = """# 数据体检\n\n> 可编辑演示\n<!-- notes: 开场。 -->\n\n## 方法\n<!-- layout: section -->\n<!-- notes: 转入方法。 -->\n\n## 先看边界\n- 记录行列数\n- 不自动填补\n<!-- notes: 说明边界。 -->\n\n## 缺失摘要\n| 字段 | 缺失 |\n|---|---:|\n| PM2.5 | 2 |\n<!-- notes: 解释表格。 -->\n\n## 复核流程\n<!-- graphic: quality-flow | 四步复核流程 -->\n<!-- notes: 说明占位合同。 -->\n"""


class SlidePlanTests(unittest.TestCase):
    def _build(self, root: Path) -> tuple[Path, Path, Path, Path]:
        source = root / "outline.md"
        source.write_text(MARKDOWN, encoding="utf-8")
        workspace = root / "artifact-workspace"
        env = os.environ.copy()
        env["HOME"] = str(Path.home())
        assert NODE and SETUP
        setup = subprocess.run([str(NODE), str(SETUP), "--workspace", str(workspace)], check=False, capture_output=True, text=True, encoding="utf-8", env=env)
        self.assertEqual(setup.returncode, 0, setup.stderr)
        pptx = root / "out" / "deck.pptx"
        plan = root / "out" / "plan.json"
        visual_qa = root / "out" / "visual-qa.json"
        render_dir = root / "out" / "rendered"
        build = subprocess.run(
            [
                str(NODE), str(BUILD_SCRIPT),
                "--workspace", str(workspace),
                "--input", str(source),
                "--allowed-root", str(root),
                "--plan-output", str(plan),
                "--pptx-output", str(pptx),
                "--render-dir", str(render_dir),
                "--qa-output", str(visual_qa),
            ],
            check=False,
            capture_output=True,
            text=True,
            encoding="utf-8",
            env=env,
        )
        self.assertEqual(build.returncode, 0, build.stderr)
        return pptx, plan, visual_qa, render_dir

    @unittest.skipUnless(ARTIFACT_RUNTIME_AVAILABLE, "set CODEX_BUNDLED_NODE and CODEX_ARTIFACT_SETUP to run artifact-tool integration tests")
    def test_artifact_tool_build_renders_and_passes_both_audits(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            pptx, plan, visual_qa, render_dir = self._build(root)
            visual = json.loads(visual_qa.read_text(encoding="utf-8"))
            self.assertEqual(visual["status"], "passed")
            self.assertEqual(visual["rendered_slide_count"], 5)
            self.assertEqual(len(list(render_dir.glob("slide-*.png"))), 5)
            audit = AUDIT.audit_presentation(pptx, plan, min_slides=5)
            self.assertEqual(audit["status"], "passed", audit["errors"])
            names = [shape.name for slide in Presentation(pptx).slides for shape in slide.shapes]
            self.assertIn("DATA_TABLE:table-4", names)
            self.assertIn("GRAPHIC_PLACEHOLDER:quality-flow", names)

    @unittest.skipUnless(ARTIFACT_RUNTIME_AVAILABLE, "set CODEX_BUNDLED_NODE and CODEX_ARTIFACT_SETUP to run artifact-tool integration tests")
    def test_structure_audit_cli_accepts_artifact_tool_deck(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            pptx, plan, _, _ = self._build(root)
            audit_path = root / "out" / "structure-audit.json"
            checked = subprocess.run(
                [sys.executable, "-B", str(AUDIT_SCRIPT), "--input", str(pptx), "--plan", str(plan), "--report-output", str(audit_path)],
                check=False,
                capture_output=True,
                text=True,
                encoding="utf-8",
            )
            self.assertEqual(checked.returncode, 0, checked.stderr)
            self.assertEqual(json.loads(audit_path.read_text(encoding="utf-8"))["status"], "passed")

    def test_audit_rejects_blank_slide(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            deck = Presentation()
            for _ in range(5):
                deck.slides.add_slide(deck.slide_layouts[6])
            deck_path = root / "blank.pptx"
            deck.save(deck_path)
            audit = AUDIT.audit_presentation(deck_path, min_slides=5)
            self.assertEqual(audit["status"], "failed")
            self.assertIn("slide_1:blank_slide", audit["errors"])


if __name__ == "__main__":
    unittest.main()
