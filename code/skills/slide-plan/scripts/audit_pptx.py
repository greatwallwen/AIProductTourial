from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation


def _shape_text(shape: Any) -> str:
    if getattr(shape, "has_text_frame", False):
        return shape.text.strip()
    if getattr(shape, "has_table", False):
        return " ".join(cell.text.strip() for row in shape.table.rows for cell in row.cells).strip()
    return ""


def audit_presentation(input_path: str | Path, plan_path: str | Path | None = None, min_slides: int = 5) -> dict[str, Any]:
    path = Path(input_path).resolve()
    if path.suffix.lower() != ".pptx" or not path.is_file():
        raise ValueError("input_must_be_existing_pptx")
    plan: dict[str, Any] | None = None
    if plan_path:
        plan = json.loads(Path(plan_path).read_text(encoding="utf-8-sig"))
    prs = Presentation(path)
    errors: list[str] = []
    slides: list[dict[str, Any]] = []
    if len(prs.slides) < min_slides:
        errors.append(f"page_count_below_minimum:{len(prs.slides)}<{min_slides}")
    if plan and plan.get("slide_count") != len(prs.slides):
        errors.append("plan_slide_count_mismatch")

    for number, slide in enumerate(prs.slides, start=1):
        titles = [
            _shape_text(shape)
            for shape in slide.shapes
            if shape.name == "SLIDE_TITLE" and _shape_text(shape)
        ]
        visible_text = [_shape_text(shape) for shape in slide.shapes if _shape_text(shape)]
        notes = slide.notes_slide.notes_text_frame.text.strip()
        placeholders = [
            shape.name
            for shape in slide.shapes
            if shape.name.startswith(("DATA_TABLE:", "GRAPHIC_PLACEHOLDER:"))
        ]
        slide_errors: list[str] = []
        if not titles:
            slide_errors.append("missing_title")
        if not visible_text:
            slide_errors.append("blank_slide")
        if not notes:
            slide_errors.append("missing_notes")
        elif "[Sources]" not in notes:
            slide_errors.append("missing_sources_block")
        expected_placeholders: list[str] = []
        if plan:
            planned = plan["slides"][number - 1]
            for placeholder in planned.get("placeholders", []):
                prefix = "DATA_TABLE:" if placeholder["type"] == "data-table" else "GRAPHIC_PLACEHOLDER:"
                expected_placeholders.append(prefix + placeholder["id"])
            for expected in expected_placeholders:
                if expected not in placeholders:
                    slide_errors.append(f"missing_placeholder:{expected}")
            if planned.get("title") not in titles:
                slide_errors.append("plan_title_mismatch")
        errors.extend(f"slide_{number}:{error}" for error in slide_errors)
        slides.append(
            {
                "slide_number": number,
                "title": titles[0] if titles else None,
                "notes_present": bool(notes),
                "sources_block_present": "[Sources]" in notes,
                "placeholders": placeholders,
                "blank": not bool(visible_text),
                "errors": slide_errors,
            }
        )
    return {
        "schema_version": "1.0",
        "status": "passed" if not errors else "failed",
        "file": str(path),
        "page_count": len(prs.slides),
        "checks": ["page_count", "titles", "speaker_notes", "sources_blocks", "placeholder_contracts", "no_blank_slides"],
        "slides": slides,
        "errors": errors,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Audit PPTX structure without PowerPoint.")
    parser.add_argument("--input", required=True)
    parser.add_argument("--plan")
    parser.add_argument("--min-slides", type=int, default=5)
    parser.add_argument("--report-output", required=True)
    args = parser.parse_args()
    try:
        result = audit_presentation(args.input, args.plan, args.min_slides)
        target = Path(args.report_output).resolve()
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(json.dumps(result, ensure_ascii=False, indent=2, sort_keys=True) + "\n", encoding="utf-8")
    except (OSError, TypeError, ValueError, json.JSONDecodeError) as error:
        print(json.dumps({"status": "blocked", "reason": str(error)}, ensure_ascii=False))
        return 2
    print(json.dumps({"status": result["status"], "page_count": result["page_count"], "report_output": str(target)}, ensure_ascii=False))
    return 0 if result["status"] == "passed" else 1


if __name__ == "__main__":
    raise SystemExit(main())
